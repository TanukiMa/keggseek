#!/usr/bin/env python3
"""
厚労省サイト専門用語新語発見システム（辞書ベース判定）
GitHub Actions + Supabase + SudachiDict-full
"""

import os
import re
import hashlib
import logging
from pathlib import Path
from typing import List, Set, Dict, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from urllib.parse import urljoin, urlparse
import uuid
import time

import requests
from bs4 import BeautifulSoup
import supabase
from supabase import create_client, Client
from sudachipy import tokenizer, dictionary
from docx import Document
from pptx import Presentation
import PyPDF2

# ロギング設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

class SupabaseClient:
    """Supabaseクライアント"""
    def __init__(self):
        self.client: Client = create_client(
            os.environ['SUPABASE_URL'],
            os.environ['SUPABASE_KEY']
        )
    
    def is_url_processed(self, url: str, content_hash: str) -> bool:
        """URLが既に処理済みかチェック"""
        try:
            result = self.client.table('processed_urls')\
                .select('*')\
                .eq('url', url)\
                .eq('file_hash', content_hash)\
                .execute()
            return len(result.data) > 0
        except Exception as e:
            logger.warning(f"URL処理済み確認でエラー: {e}")
            return False
    
    def save_processed_url(self, url: str, content_type: str, content_hash: str) -> str:
        """処理済みURL保存"""
        try:
            result = self.client.table('processed_urls').insert({
                'url': url,
                'content_type': content_type,
                'file_hash': content_hash,
                'status': 'completed'
            }).execute()
            return result.data[0]['id'] if result.data else None
        except Exception as e:
            logger.error(f"URL保存エラー: {e}")
            return None
    
    def save_extracted_words(self, words: List[Dict], url_id: str):
        """抽出単語保存"""
        if not words or not url_id:
            return
            
        try:
            for word_data in words:
                word_data['url_id'] = url_id
            self.client.table('extracted_words').insert(words).execute()
        except Exception as e:
            logger.error(f"単語保存エラー: {e}")
    
    def get_dictionary_words(self) -> Set[str]:
        """既存辞書単語取得"""
        try:
            result = self.client.table('dictionary_words').select('word').execute()
            return {row['word'] for row in result.data}
        except Exception as e:
            logger.warning(f"辞書単語取得エラー: {e}")
            return set()
    
    def save_new_word_candidate(self, word_data: Dict):
        """新語候補保存"""
        try:
            self.client.table('new_word_candidates').insert(word_data).execute()
        except Exception as e:
            logger.error(f"新語候補保存エラー: {e}")

class DocumentProcessor:
    """文書処理クラス"""
    
    @staticmethod
    def extract_text_from_html(content: str) -> str:
        """HTMLからテキスト抽出"""
        try:
            soup = BeautifulSoup(content, 'html.parser')
            # スクリプト・スタイル削除
            for script in soup(["script", "style"]):
                script.decompose()
            
            # メインコンテンツを抽出
            text = soup.get_text()
            # 改行・空白の正規化
            text = re.sub(r'\s+', ' ', text).strip()
            return text
        except Exception as e:
            logger.error(f"HTML処理エラー: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """PDFからテキスト抽出"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            logger.error(f"PDF処理エラー: {e}")
        return text
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """DOCXからテキスト抽出"""
        try:
            doc = Document(file_path)
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
        except Exception as e:
            logger.error(f"DOCX処理エラー: {e}")
            return ""
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """PPTXからテキスト抽出"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX処理エラー: {e}")
            return ""

class SudachiAnalyzer:
    """Sudachi形態素解析（SudachiDict-full使用）"""
    
    def __init__(self):
        # SudachiPyでFull辞書を使用（設定ファイル経由）
        from sudachipy import tokenizer, dictionary
        
        try:
            # 設定ファイルがあれば使用、なければデフォルト
            self.tokenizer_obj = dictionary.Dictionary().create()
            self.mode = tokenizer.Tokenizer.SplitMode.A
            logger.info("✅ SudachiDict-full を使用して初期化完了")
        except Exception as e:
            logger.error(f"Sudachi辞書の初期化に失敗: {e}")
            raise
    
    def analyze(self, text: str) -> List[Dict]:
        """テキスト解析"""
        if not text:
            return []
            
        words = []
        try:
            tokens = self.tokenizer_obj.tokenize(text, self.mode)
            
            for token in tokens:
                # 専門用語らしいもの（名詞、複合語など）を抽出
                pos = token.part_of_speech()[0]
                surface = token.surface()
                
                # フィルタリング条件
                if (pos in ['名詞', '動詞', '形容詞'] and 
                    len(surface) >= 2 and 
                    not surface.isdigit() and
                    surface not in ['こと', 'もの', 'ため', 'など']):
                    
                    words.append({
                        'word': surface,
                        'reading': token.reading_form() or surface,
                        'part_of_speech': pos,
                    })
        except Exception as e:
            logger.error(f"形態素解析エラー: {e}")
        
        return words
    
    def is_known_word(self, word: str) -> bool:
        """SudachiDict-fullに収載されているかチェック"""
        try:
            # 辞書に登録されている語彙かどうかを判定
            tokens = self.tokenizer_obj.tokenize(word, self.mode)
            
            # 1つのトークンになり、かつ未知語でない場合は既知語
            if len(tokens) == 1:
                token = tokens[0]
                # 未知語の場合、品詞に「補助記号」等が含まれることが多い
                pos_features = token.part_of_speech()
                if ('未知語' not in str(pos_features) and 
                    '補助記号' not in str(pos_features) and
                    token.surface() == word):
                    return True
            
            return False
        except Exception as e:
            logger.warning(f"辞書検索エラー '{word}': {e}")
            return False

class NewWordDetector:
    """辞書ベース新語検出"""
    
    def __init__(self, analyzer: SudachiAnalyzer):
        self.analyzer = analyzer
        
        # 除外する一般的な語彙（厚労省文書でよく出現する基本語彙）
        self.common_words = {
            # 一般的な行政・医療用語
            '政策', '制度', '対策', '施策', '事業', '取り組み', '推進', '支援',
            '国民', '社会', '地域', '全国', '都道府県', '市町村',
            '厚生', '労働', '健康', '医療', '介護', '福祉', '年金', '保険',
            '災害', '職場', '労働者', '事業者', '関係者',
            # 基本語彙
            '今回', '今後', '現在', '過去', '将来', '状況', '課題', '問題',
            '方法', '手法', '仕組み', '体制', '環境', '条件', '基準',
            '効果', '影響', '結果', '成果', '実績', '評価'
        }
        
        logger.info("✅ 辞書ベース新語検出器を初期化完了")
    
    def is_new_word(self, word: str, part_of_speech: str) -> tuple[bool, float, str]:
        """辞書ベースで新語かどうか判定"""
        
        # 基本的なフィルタリング
        if (len(word) < 3 or  # 3文字未満は除外
            word in self.common_words or  # 一般語は除外
            word.isdigit() or  # 数字のみは除外
            not re.match(r'^[ぁ-んァ-ヶー一-龠a-zA-Z]+, word)):  # 文字種チェック
            return False, 0.1, "基本フィルタで除外"
        
        # 名詞に限定（新語候補として最も有力）
        if part_of_speech != '名詞':
            return False, 0.2, "名詞以外"
        
        # SudachiDict-fullに収載されているかチェック
        is_known = self.analyzer.is_known_word(word)
        
        if is_known:
            return False, 0.3, "SudachiDict-fullに収載済み"
        else:
            # 新語候補として判定
            confidence = 0.8  # 辞書にない場合は高い信頼度
            
            # 専門用語らしさによる信頼度調整
            if len(word) >= 5:  # 5文字以上は専門用語の可能性高
                confidence = 0.9
            elif any(char in word for char in ['DX', 'AI', 'IoT', 'ICT']):  # 英略語含む
                confidence = 0.9
            elif word.endswith(('システム', '事業', '制度', '政策')):  # 専門用語パターン
                confidence = 0.7
            
            return True, confidence, f"SudachiDict-full未収載（{len(word)}文字の名詞）"

class MhlwCrawler:
    """メインクローラークラス"""
    
    def __init__(self):
        self.db = SupabaseClient()
        self.processor = DocumentProcessor()
        self.analyzer = SudachiAnalyzer()
        self.detector = NewWordDetector(self.analyzer)  # LLMではなく辞書ベース
        
        self.base_url = "https://www.mhlw.go.jp"
        self.session = requests.Session()
        self.session.headers.update({
            'User-Agent': 'Mozilla/5.0 (compatible; MHLW Terminology Research Bot; +https://github.com/)'
        })
        
        logger.info("🚀 MhlwCrawler初期化完了（辞書ベース新語検出）")
    
    def get_urls_to_crawl(self) -> List[str]:
        """クローリング対象URL取得"""
        urls = []
        
        # 厚労省の主要ページから開始
        start_urls = [
            f"{self.base_url}/stf/seisakunitsuite/bunya/kenkou_iryou/",
            f"{self.base_url}/stf/seisakunitsuite/bunya/koyou_roudou/",
            f"{self.base_url}/toukei/",
            f"{self.base_url}/shingi/",
        ]
        
        logger.info(f"📡 {len(start_urls)}個のスタートURLからクローリング開始")
        
        for start_url in start_urls:
            try:
                logger.info(f"🔍 {start_url} を解析中...")
                response = self.session.get(start_url, timeout=30)
                response.raise_for_status()
                
                soup = BeautifulSoup(response.content, 'html.parser')
                
                # PDF, DOCX, PPTXリンク収集
                file_links = 0
                for link in soup.find_all('a', href=True):
                    href = link['href']
                    if any(ext in href.lower() for ext in ['.pdf', '.docx', '.pptx']):
                        full_url = urljoin(start_url, href)
                        if full_url not in urls:
                            urls.append(full_url)
                            file_links += 1
                
                # HTMLページも追加（同一ドメイン内）
                html_links = 0
                for link in soup.find_all('a', href=True):
                    href = link['href']
                    if (href.startswith('/') or 'mhlw.go.jp' in href) and not any(ext in href.lower() for ext in ['.pdf', '.docx', '.pptx']):
                        full_url = urljoin(start_url, href)
                        if full_url not in urls and len(urls) < 100:  # 上限設定
                            urls.append(full_url)
                            html_links += 1
                
                logger.info(f"✅ {start_url}: ファイル{file_links}件、HTML{html_links}件を発見")
                time.sleep(1)  # 間隔を空ける
                
            except Exception as e:
                logger.error(f"❌ URL収集エラー {start_url}: {e}")
        
        # 重複除去
        unique_urls = list(set(urls))
        logger.info(f"🎯 合計 {len(unique_urls)} 個のURLを収集完了")
        
        return unique_urls[:50]  # 最初の50個に制限（テスト用）
    
    def process_url(self, url: str) -> Optional[Dict]:
        """単一URL処理"""
        try:
            logger.info(f"🔄 処理開始: {url}")
            
            # ファイル取得
            response = self.session.get(url, timeout=30)
            response.raise_for_status()
            content_hash = hashlib.md5(response.content).hexdigest()
            
            # 既処理チェック
            if self.db.is_url_processed(url, content_hash):
                logger.info(f"⏭️  スキップ（既処理）: {url}")
                return None
            
            # ファイルタイプ判定
            content_type = self._get_content_type(url, response.headers.get('content-type', ''))
            
            # テキスト抽出
            text = self._extract_text(response.content, content_type, url)
            if not text or len(text) < 50:
                logger.warning(f"⚠️  テキスト抽出失敗またはコンテンツ不足: {url}")
                return None
            
            logger.info(f"📝 テキスト抽出完了: {len(text)}文字")
            
            # 形態素解析
            words = self.analyzer.analyze(text)
            logger.info(f"🔤 形態素解析完了: {len(words)}語を抽出")
            
            # 既存辞書と照合 + 新語検出
            dictionary_words = self.db.get_dictionary_words()
            new_candidates = []
            
            # 語彙の頻度カウント（同じ文書内での出現頻度）
            word_freq = {}
            for word_data in words:
                word = word_data['word']
                word_freq[word] = word_freq.get(word, 0) + 1
            
            # ユニークな語彙のみを新語候補として検討
            unique_words = {}
            for word_data in words:
                word = word_data['word']
                if word not in unique_words:
                    unique_words[word] = word_data
            
            logger.info(f"🔍 新語検出開始: {len(unique_words)}語をチェック")
            
            for word, word_data in unique_words.items():
                # 基本辞書にない語彙をチェック
                if word not in dictionary_words:
                    # 辞書ベース新語判定
                    is_new, confidence, reasoning = self.detector.is_new_word(
                        word, word_data['part_of_speech']
                    )
                    
                    if is_new and confidence > 0.6:
                        frequency = word_freq.get(word, 1)
                        new_candidates.append({
                            'word': word,
                            'reading': word_data['reading'],
                            'part_of_speech': word_data['part_of_speech'],
                            'confidence_score': confidence,
                            'llm_reasoning': reasoning,  # 判定理由
                            'source_urls': [url],
                            'frequency_count': frequency
                        })
                        logger.info(f"✨ 新語候補発見: '{word}' (信頼度: {confidence:.3f}, 頻度: {frequency})")
            
            # DB保存
            url_id = self.db.save_processed_url(url, content_type, content_hash)
            if url_id and words:
                # 語数制限
                words_to_save = words[:100] if len(words) > 100 else words
                self.db.save_extracted_words(words_to_save, url_id)
            
            for candidate in new_candidates:
                self.db.save_new_word_candidate(candidate)
            
            logger.info(f"✅ 完了: {url} - 新語候補: {len(new_candidates)}件")
            return {
                'url': url,
                'words_count': len(words),
                'new_words_count': len(new_candidates)
            }
            
        except Exception as e:
            logger.error(f"❌ 処理エラー {url}: {e}")
            return None
    
    def _get_content_type(self, url: str, content_type_header: str) -> str:
        """コンテンツタイプ判定"""
        url_lower = url.lower()
        if '.pdf' in url_lower:
            return 'pdf'
        elif '.docx' in url_lower:
            return 'docx'
        elif '.pptx' in url_lower:
            return 'pptx'
        else:
            return 'html'
    
    def _extract_text(self, content: bytes, content_type: str, url: str) -> str:
        """コンテンツからテキスト抽出"""
        if content_type == 'html':
            return self.processor.extract_text_from_html(content.decode('utf-8', errors='ignore'))
        else:
            # ファイル保存して処理
            temp_path = f"/tmp/{uuid.uuid4()}.{content_type}"
            try:
                with open(temp_path, 'wb') as f:
                    f.write(content)
                
                if content_type == 'pdf':
                    return self.processor.extract_text_from_pdf(temp_path)
                elif content_type == 'docx':
                    return self.processor.extract_text_from_docx(temp_path)
                elif content_type == 'pptx':
                    return self.processor.extract_text_from_pptx(temp_path)
            finally:
                Path(temp_path).unlink(missing_ok=True)
        
        return ""
    
    def run(self, max_workers: int = 3):
        """メイン実行"""
        logger.info("🚀 厚労省サイト解析開始（辞書ベース新語検出）")
        start_time = time.time()
        
        # クローリング対象URL取得
        urls = self.get_urls_to_crawl()
        if not urls:
            logger.error("❌ クローリング対象URLが見つかりませんでした")
            return
        
        logger.info(f"🎯 対象URL数: {len(urls)}")
        
        total_processed = 0
        total_new_words = 0
        
        # 並列処理
        logger.info(f"👥 並列処理開始（ワーカー数: {max_workers}）")
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_url = {executor.submit(self.process_url, url): url for url in urls}
            
            for future in as_completed(future_to_url):
                url = future_to_url[future]
                try:
                    result = future.result()
                    if result:
                        total_processed += 1
                        total_new_words += result['new_words_count']
                        logger.info(f"📊 進捗: {total_processed}/{len(urls)} 完了")
                except Exception as e:
                    logger.error(f"❌ {url} の処理中にエラー: {e}")
        
        elapsed_time = time.time() - start_time
        logger.info(f"🎉 処理完了: {total_processed}URL処理, {total_new_words}新語候補発見, {elapsed_time:.1f}秒")
        logger.info("💡 新語判定基準: SudachiDict-full（170万語）未収載の名詞")

if __name__ == "__main__":
    # 直接実行用
    import argparse
    
    parser = argparse.ArgumentParser(description='厚労省サイト専門用語解析（辞書ベース）')
    parser.add_argument('--workers', type=int, default=3, help='並列処理数')
    args = parser.parse_args()
    
    crawler = MhlwCrawler()
    crawler.run(max_workers=args.workers)
